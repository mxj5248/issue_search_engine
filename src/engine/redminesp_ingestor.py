from elasticsearch import Elasticsearch
from elasticsearch import helpers
from redminelib import Redmine
from pptx import Presentation
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO
from util.querystore import *
from datetime import datetime
from redminelib.exceptions import ResourceNotFoundError
from urllib import parse
from abc import *
import pandas as pd
import datetime
import hashlib
import pymysql 
import docx
import os
import pickle
import time

# 이슈를 표현하는 class 입니다.
class ProductIssue(object):
  """
    Represents semantic data for a single product post
  """
  def __init__(self, id, container_id, subject, description, created_on):
        self.id = id
        self.description = description
        self.container_id = container_id
        self.subject = subject
        self.created_on= created_on

REDMINE_DB_PORT = int(os.getenv("REDMINE_DB_PORT"))
REDMINE_DB_DATABASE = os.getenv("REDMINE_DB_DATABASE")
REDMINE_DB_USER = os.getenv("REDMINE_DB_USER")
REDMINE_DB_PASSWORD = os.getenv("REDMINE_DB_PASSWORD")
REDMINE_DB_HOST = os.getenv("REDMINE_DB_HOST")

def get_conn():
    conn = pymysql.connect(host=REDMINE_DB_HOST,port=REDMINE_DB_PORT,user=REDMINE_DB_USER,passwd=REDMINE_DB_PASSWORD,db=REDMINE_DB_DATABASE,charset='utf8')
    return conn

#// Pool에 Connection 반환 ( 무조건 해 주어야 한다. )
def release(conn):
    conn.close()

def json_default(value): 
    if isinstance(value, datetime.date): 
        return value.strftime('%Y-%m-%d') 
    raise TypeError('not JSON serializable')

# Custom handlers for marshalling python object into JSON 
def json_field_handler(x):
    if isinstance(x, datetime):
        return x.isoformat()
    raise TypeError("Unable to parse json field")

# 엘라스틱서치에서 사용될 문서의 고유 아이디를 생성합니다.
def getUniqueIndexId(pri):
    return hashlib.sha1(pri.encode('utf-8')).hexdigest()

def get_ppt_text(filename):
    eachfile = filename
    prs = Presentation(eachfile)
    contents = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                contents.append(shape.text)

    contents = "".join(contents)
    if len(contents) > 2000 :
        contents = contents[:2000]
        return contents
    else: 
        return contents

def get_pdf_text(filename):
    try:
        rsrcmgr = PDFResourceManager()
        retstr = StringIO()
        codec = "utf-8"
        laparams = LAParams()
        device = TextConverter(
            rsrcmgr, 
            retstr, 
            codec=codec, 
            laparams=laparams
        )
        fp = open(filename, 'rb')
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        password = ""
        maxpages = 0
        caching = True
        pagenos = set()

        pages = PDFPage.get_pages(
            fp, 
            pagenos, 
            maxpages=maxpages, 
            password=password, 
            caching=caching, 
            check_extractable=True
        )
        for page in pages:
            interpreter.process_page(page)

        text = retstr.getvalue()
        fp.close()
        device.close()
        retstr.close()
        if len(text) > 2000 :
            text = text[:2000]
            return text
        elif len(text) == 0:
            return ""
        else: 
            return text
    except Exception as e:
        print(e)
        text = ""

def get_doc_text(filename): 
    try:
        doc = docx.Document(filename)
        fullText = ''
        for para in doc.paragraphs:
            fullText+=(para.text)
        if len(fullText) > 2000 :
            fullText =fullText[:2000]
            return fullText
        else: 
            return fullText
    
    except Exception as e:
        print(e)

def get_txt_text(filename):
    try:
        with open(filename, "r", encoding='utf-8') as f:
            string = ''
            for line in f:
                string += line.strip()
            if len(string) > 2000 :
                string =string[:2000]
                return string
            else: 
                return string
    except Exception as e:
        print(e)
        pass

def get_excel_text(filename):
    try:
        xls = pd.ExcelFile(filename)
        string = ""
        for sheet in xls.sheet_names:
            df = pd.read_excel(filename, sheet_name=sheet)
            df = df.fillna("")
            for i in df.values:
                string += " "
                for j in i:
                    string +=str(j)
        if len(string) > 2000 :
            string =string[:2000]
            return string
        else: 
            return string
    except  Exception as e:
        print(e)
        pass




def encode_text(text):
    enc = parse.quote(text)
    return enc

def save_pickle(df,name):
    f_name = name + '.pickle'
    with open(f_name, 'wb') as f:
        pickle.dump(df.to_dict(), f, pickle.HIGHEST_PROTOCOL)

# SQL 에 연결하여 이슈 테이블을 추출하여 ProductIssue list 로 돌려주는 함수입니다
def getIssues():
    conn = get_conn()
    cursor = conn.cursor(pymysql.cursors.DictCursor)
    query = "SELECT id, container_id, container_type, filename, disk_filename, filesize, created_on FROM redmine.attachments a WHERE a.created_on >= '2021-01-01' "
    cursor.execute(query)
    data = cursor.fetchall()
    result = pd.DataFrame(data,columns=['id', 'container_id', 'container_type', 'filename', 'disk_filename', 'filesize', 'created_on'])
    result = result[result['container_type']=="Issue"]
    
    def get_r_contents(rid):
        conn = get_conn()
        cursor = conn.cursor(pymysql.cursors.DictCursor)
        r_query = "SELECT id, description, project_id, subject, created_on FROM redmine.issues i WHERE i.id = " + str(rid)
        cursor.execute(r_query)
        r_data = cursor.fetchall()
        r_result = pd.DataFrame(r_data,columns=['id', 'description', 'project_id', 'subject', 'created_on'])
        release(conn)
        return r_result 

    redmine = Redmine("https://redmine.netand.co.kr",username='mxj5248',password='Tyler1031!')
#################################################################################################################################################################
    # start = time.time()  # 시작 시간 저장
    # powers = []
    # des_powers = []
    # rdes_powers = []
    # rsubj_powers = []
    # for row, value in result.iterrows():
    #     if value['filename'][-3:] in ['ptx']:
    #         try:
    #             print("Issue id {} found. file: {} , created on: {}".format(value['id'],value['filename'],value['created_on']))
    #             file_id = str(value['id'])
    #             file_name = value['filename']
    #             prefix_url = "https://redmine.netand.co.kr/attachments/download/"
    #             enc_name = encode_text(file_name)
    #             url = prefix_url + file_id + '/' + enc_name
    #             redmine.download(url,savepath='/datastore',filename = file_name)
                
    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_powers.append(r_result['description'])
    #             rsubj_powers.append(r_result['subject'])
                
    #             content = get_ppt_text('/datastore/'+file_name) 
    #             des_powers.append(content)
    #             powers.append(value['id'])

    #             # r_result = get_r_contents(value['container_id'])
    #             # rdes_powers.append(r_result['description'])
    #             # rsubj_powers.append(r_result['subject'])

    #         except Exception as e:
    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_powers.append(r_result['description'])
    #             rsubj_powers.append(r_result['subject'])
    #             pass

    # result_p = result[result['id'].isin(powers)]
    # result_p['description'] = des_powers
    # result_p['rdescription'] = rdes_powers
    # result_p['rsubjet'] = rsubj_powers
    # save_pickle(result_p,'ppt')
    # print("Saved Pickle")
    # print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간
    # print('------------------------------next type-------------------------------------------')
    
    # start = time.time()  # 시작 시간 저장
    # pdfs = []
    # des_pdfs = []
    # rdes_pdfs = []
    # rsubj_pdfs = []
    # for row, value in result.iterrows():
    #     if value['filename'][-3:] in ['pdf']:
    #         try:
    #             print("Issue id {} found. file: {} , created on: {}".format(value['id'],value['filename'],value['created_on']))
    #             file_id = str(value['id'])
    #             file_name = value['filename']
    #             prefix_url = "https://redmine.netand.co.kr/attachments/download/"
    #             enc_name = encode_text(file_name)
    #             url = prefix_url + file_id + '/' + enc_name
    #             redmine.download(url,savepath='/datastore',filename = file_name)

    #             pdfs.append(value['id'])
    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_pdfs.append(r_result['description'])
    #             rsubj_pdfs.append(r_result['subject'])

    #             content = get_pdf_text('/datastore/'+file_name)
    #             des_pdfs.append(content)
            
    #         except Exception as e:
    #             des_pdfs.append("파일을 확인해 주세요.")
    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_pdfs.append(r_result['description'])
    #             rsubj_pdfs.append(r_result['subject'])
    #             rdes_pdfs.append("파일을 확인해 주세요.")
    #             rsubj_pdfs.append("파일을 확인해 주세요.")
    #             pass
    # print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간

    # result_pd = result[result['id'].isin(pdfs)]
    # print(len(result_pd))

    # result_pd['description'] =  pd.Series(des_pdfs)
    # result_pd['rdescription'] = rdes_pdfs
    # result_pd['rsubjet'] = rsubj_pdfs
    # save_pickle(result_pd,'pdfs')
    # print("Saved Pickle")
    # print('------------------------------next type-------------------------------------------')

    start = time.time()  # 시작 시간 저장
    docs = []
    des_docs = []
    rdes_docs = []
    rsubj_docs = []
    for row, value in result.iterrows():
        if value['filename'][-4:] in ['docx']:
            try:
                print("Issue id {} found. file: {} , created on: {}".format(value['id'],value['filename'],value['created_on']))
                file_id = str(value['id'])
                file_name = value['filename']
                prefix_url = "https://redmine.netand.co.kr/attachments/download/"
                enc_name = encode_text(file_name)
                url = prefix_url + file_id + '/' + enc_name
                redmine.download(url,savepath='/datastore',filename = file_name)
                
                docs.append(value['id'])
            
                r_result = get_r_contents(int(value['container_id']))
                rdes_docs.append(r_result['description'])
                rsubj_docs.append(r_result['subject'])
                
                content = get_doc_text('/datastore/'+file_name)
                des_docs.append(content)

            except Exception as e:
                des_docs.append("파일을 확인해 주세요.")
                r_result = get_r_contents(int(value['container_id']))
                rdes_docs.append(r_result['description'])
                rsubj_docs.append(r_result['subject'])
                rdes_docs.append("파일을 확인해 주세요.")
                rsubj_docs.append("파일을 확인해 주세요.")
                pass

    result_doc = result[result['id'].isin(docs)]
    print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간
    print(len(result_doc))
    result_doc['description'] = des_docs
    result_doc['rdescription'] = rdes_docs
    result_doc['rsubjet'] = rsubj_docs
    save_pickle(result_doc,'docs')
    print("Saved Pickle")
    print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간
    print('------------------------------next type-------------------------------------------')

    start = time.time()  # 시작 시간 저장
    txts = []
    des_txts = []
    rdes_txts = []
    rsubj_txts = []
    for row, value in result.iterrows():
        if value['filename'][-3:] in ['txt']:
            try:
                print("Issue id {} found. file: {} , created on: {}".format(value['id'],value['filename'],value['created_on']))
                file_id = str(value['id'])
                file_name = value['filename']
                prefix_url = "https://redmine.netand.co.kr/attachments/download/"
                enc_name = encode_text(file_name)
                url = prefix_url + file_id + '/' + enc_name
                redmine.download(url,savepath='/datastore',filename = file_name)
                content = get_txt_text('/datastore/'+file_name)
                txts.append(value['id'])

                r_result = get_r_contents(int(value['container_id']))
                rdes_txts.append(r_result['description'])
                rsubj_txts.append(r_result['subject'])
                des_txts.append(content)

            except Exception as e:
                des_txts.append("파일을 확인해 주세요.")
                r_result = get_r_contents(int(value['container_id']))
                rdes_txts.append(r_result['description'])
                rsubj_txts.append(r_result['subject'])
                rdes_txts.append("파일을 확인해 주세요.")
                rsubj_txts.append("파일을 확인해 주세요.")
                pass

    result_txt = result[result['id'].isin(txts)]
    print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간

    result_txt['description'] = des_txts
    result_txt['rdescription'] = rdes_txts
    result_txt['rsubjet'] = rsubj_txts

    save_pickle(result_txt,'txts')
    print("Saved Pickle")
    print('------------------------------next type-------------------------------------------')

    # start = time.time()  # 시작 시간 저장
    # excels = []
    # des_excels = []
    # rdes_excels = []
    # rsubj_excels = []

    # for row, value in result.iterrows():
    #     if value['filename'][-3:] in ['xls','lsx','csv']:
    #         try: 
    #             print("Issue id {} found. file: {} , created on: {}".format(value['id'],value['filename'],value['created_on']))
    #             file_id = str(value['id'])
    #             file_name = value['filename']
    #             prefix_url = "https://redmine.netand.co.kr/attachments/download/"
    #             enc_name = encode_text(file_name)
    #             url = prefix_url + file_id + '/' + enc_name
    #             redmine.download(url,savepath='/datastore',filename = file_name)
    #             content = get_excel_text('/datastore/'+file_name)
    #             des_excels.append(content)
    #             excels.append(value['id'])

    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_excels.append(r_result['description'])
    #             rsubj_excels.append(r_result['subject'])

    #         except Exception as e:
    #             des_excels.append("파일을 확인해 주세요.")
    #             r_result = get_r_contents(int(value['container_id']))
    #             rdes_excels.append(r_result['description'])
    #             rsubj_excels.append(r_result['subject'])
    #             rdes_excels.append("파일을 확인해 주세요.")
    #             rsubj_excels.append("파일을 확인해 주세요.")
    #             pass
    
    # result_e = result[result['id'].isin(excels)]
    # print("time :", time.time() - start)  # 현재시각 - 시작시간 = 실행 시간
    # result_e['description'] = des_excels
    # result_e['rdescription'] = rdes_excels
    # result_e['rsubjet'] = rsubj_excels
    # save_pickle(result_e,'excels')
    # print("Saved Pickle")
    #################################################################################################################
    
    # result = pd.concat([result_p,result_d,result_t,result_pd])
    # result = pd.concat([result_p,result_d,result_t,result_pd,result_e])
    # result = result_p[['id','container_id','filename','description','created_on','rdescription','rsubjet']]
    # result.rename(columns={'filename':'subject'},inplace=True)

    release(conn)
    # print(len(result))
    # return result

def issueToElasticSearch(df):
    url = "http://elasticsearch"
    port = "9200"

    def get_conn():
        es = Elasticsearch(f'{url}:{port}')
        return es
    es = get_conn()
    data = [
    {
        "_index": "idx_redminesp",
        "_type": "_doc",
        "_id": getUniqueIndexId(x[2]),
        "_source": {
            "id": x[0],
            "container_id" : x[1],
            "description": x[3],
            "subject": x[2],
            "created_on": x[4],
            "red_description": x[5],
            "red_subject": x[6],
            }
    }
        for x in zip(df['id'], df['container_id'],df['subject'], df['description'], df['created_on'],df['rdescription'],df['rsubjet'])
    ]
    helpers.bulk(es, data)


def begin_rspingestor():
    # print("")
    try:
        # p = getIssues()   
        df = pd.DataFrame()
        with open('ppt.pickle', 'rb') as ffff:
            dataff = pickle.load(ffff)
            df_ppts = pd.DataFrame(dataff)
            df = pd.concat([df,df_ppts])
        with open('docs.pickle', 'rb') as f:
            data = pickle.load(f)
            df_docs = pd.DataFrame(data)
            df = pd.concat([df,df_docs])
        with open('pdfs.pickle', 'rb') as ff:
            dataf = pickle.load(ff)
            df_pdfs = pd.DataFrame(dataf)
            df = pd.concat([df,df_pdfs])
        with open('txts.pickle', 'rb') as fff:
            datat = pickle.load(fff)
            df_txts = pd.DataFrame(datat)
            df = pd.concat([df,df_txts])
        # with open('excels.pickle', 'rb') as e:
        #     datae = pickle.load(e)
        #     df_e = pd.DataFrame(datae)
        #     df = pd.concat([df,df_e])


        p = df[['id','container_id','filename','description','created_on','rdescription','rsubjet']]
        p.rename(columns={'filename':'subject'},inplace=True)
        time.sleep(10)
        for idx in range(int(len(p)/30)):
            slice_data = p.iloc[idx*30:(idx+1)*30,:]
            print(slice_data)
            issueToElasticSearch(slice_data)
            time.sleep(30)
    except Exception as e:
        print(e)
        pass
begin_rspingestor()
